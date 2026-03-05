from __future__ import annotations

import io

import pytest

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.shapes.graphfrm import GraphicFrame


def test_graphic_frame_chart_access(parent) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = ("A", "B")
    chart_data.add_series("Series 1", (1, 2))
    chart_rid = parent.part.add_chart_part(XL_CHART_TYPE.BAR_CLUSTERED, chart_data)
    chart_xml = (
        b'<p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        b'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        b'xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
        b'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        b'<p:nvGraphicFramePr><p:cNvPr id="42" name="Chart"/></p:nvGraphicFramePr>'
        b'<a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart">'
        b'<c:chart r:id="%b"/>'
        b"</a:graphicData></a:graphic>"
        b"</p:graphicFrame>"
    ) % chart_rid.encode("utf-8")

    graphic_frame = GraphicFrame(
        parse_xml(chart_xml),
        parent,
    )

    assert graphic_frame.has_chart is True
    assert graphic_frame.has_table is False
    assert graphic_frame.shape_type == MSO_SHAPE_TYPE.CHART
    assert graphic_frame.chart is not None
    assert graphic_frame.chart_part is parent.part.related_part(chart_rid)


def test_graphic_frame_chart_raises_when_missing(parent) -> None:
    graphic_frame = GraphicFrame(
        parse_xml(
            b"""
            <p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvGraphicFramePr><p:cNvPr id="42" name="Not Chart"/></p:nvGraphicFramePr>
              <a:graphic>
                <a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/ole">
                  <p:oleObj xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>
                </a:graphicData>
              </a:graphic>
            </p:graphicFrame>
            """
        ),
        parent,
    )

    with pytest.raises(ValueError, match="shape does not contain a chart"):
        _ = graphic_frame.chart


def test_graphic_frame_table_shape_type(parent) -> None:
    graphic_frame = GraphicFrame(
        parse_xml(
            b"""
            <p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvGraphicFramePr><p:cNvPr id="42" name="Table"/></p:nvGraphicFramePr>
              <a:graphic>
                <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table"/>
              </a:graphic>
            </p:graphicFrame>
            """
        ),
        parent,
    )

    assert graphic_frame.has_table is True
    assert graphic_frame.shape_type == MSO_SHAPE_TYPE.TABLE


def test_graphic_frame_ole_format(parent) -> None:
    ole_rid = parent.part.add_embedded_ole_object_part("Excel.Sheet.12", io.BytesIO(b"xlsx"))
    ole_xml = (
        b'<p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        b'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        b'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        b'<p:nvGraphicFramePr><p:cNvPr id="42" name="OLE"/></p:nvGraphicFramePr>'
        b'<a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/ole">'
        b'<p:oleObj r:id="%b" progId="Excel.Sheet.12" showAsIcon="1"><p:embed/></p:oleObj>'
        b"</a:graphicData></a:graphic>"
        b"</p:graphicFrame>"
    ) % ole_rid.encode("utf-8")

    graphic_frame = GraphicFrame(
        parse_xml(ole_xml),
        parent,
    )

    assert graphic_frame.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT
    ole_format = graphic_frame.ole_format
    assert ole_format.blob == b"xlsx"
    assert ole_format.prog_id == "Excel.Sheet.12"
    assert ole_format.show_as_icon is True


def test_graphic_frame_ole_format_raises_on_non_ole(parent) -> None:
    graphic_frame = GraphicFrame(
        parse_xml(
            b"""
            <p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvGraphicFramePr><p:cNvPr id="42" name="Chart"/></p:nvGraphicFramePr>
              <a:graphic>
                <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart"/>
              </a:graphic>
            </p:graphicFrame>
            """
        ),
        parent,
    )

    with pytest.raises(ValueError, match="not an OLE-object shape"):
        _ = graphic_frame.ole_format


def test_graphic_frame_shadow_raises(parent) -> None:
    graphic_frame = GraphicFrame(
        parse_xml(
            b"""
            <p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvGraphicFramePr><p:cNvPr id="42" name="Chart"/></p:nvGraphicFramePr>
              <a:graphic>
                <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/chart"/>
              </a:graphic>
            </p:graphicFrame>
            """
        ),
        parent,
    )

    with pytest.raises(NotImplementedError, match="shadow property on GraphicFrame"):
        _ = graphic_frame.shadow
