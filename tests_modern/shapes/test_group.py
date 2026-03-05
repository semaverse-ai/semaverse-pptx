from __future__ import annotations

import pytest

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape

from .factories import make_group, make_sp_snippet


def test_group_shape_properties(parent) -> None:
    group = GroupShape(
        make_group(
            shape_id=42,
            name="Group 1",
            with_xfrm=True,
            children=[
                make_sp_snippet(
                    shape_id=43,
                    name="Shape 1",
                    x=10,
                    y=20,
                    cx=30,
                    cy=40,
                    prst="rect",
                    with_text_body=True,
                )
            ],
        ),
        parent,
    )

    assert group.shape_type == MSO_SHAPE_TYPE.GROUP
    assert group.has_text_frame is False
    assert len(group.shapes) == 1


def test_group_shape_click_action_raises(parent) -> None:
    group = GroupShape(make_group(shape_id=42, name="Group 1"), parent)

    with pytest.raises(TypeError, match="group shape cannot have a click action"):
        _ = group.click_action


def test_group_shape_shadow(parent) -> None:
    group = GroupShape(make_group(shape_id=42, name="Group 1"), parent)

    assert group.shadow is not None
