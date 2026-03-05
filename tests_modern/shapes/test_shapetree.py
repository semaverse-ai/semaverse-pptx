from __future__ import annotations

import io
from pathlib import Path

from syrupy.assertion import SnapshotAssertion

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, PP_PLACEHOLDER, PROG_ID
from pptx.oxml import parse_xml
from pptx.shapes.connector import Connector
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.shapes.placeholder import (
    ChartPlaceholder,
    LayoutPlaceholder,
    MasterPlaceholder,
    NotesSlidePlaceholder,
    PicturePlaceholder,
    SlidePlaceholder,
    TablePlaceholder,
)
from pptx.shapes.shapetree import (
    BaseShapeFactory,
    LayoutPlaceholders,
    LayoutShapes,
    MasterPlaceholders,
    MasterShapes,
    NotesSlidePlaceholders,
    NotesSlideShapes,
    SlidePlaceholders,
    SlideShapeFactory,
    SlideShapes,
)


def test_slide_shapes_collection_protocol(parent) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    assert len(shapes) == 0
    assert list(shapes) == []


def test_slide_shapes_add_shape(parent, snapshot: SnapshotAssertion) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 100, 200, 300, 400)

    assert shape.shape_id == 2
    assert shape.name == "Rounded Rectangle 1"
    assert len(shapes) == 1
    assert snapshot == parent.part._element.cSld.spTree.xml


def test_slide_shapes_add_textbox(parent, snapshot: SnapshotAssertion) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    shape = shapes.add_textbox(100, 200, 300, 400)

    assert shape.shape_id == 2
    assert shape.name == "TextBox 1"
    assert shape.has_text_frame is True
    assert snapshot == parent.part._element.cSld.spTree.xml


def test_slide_shapes_add_connector(parent, snapshot: SnapshotAssertion) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    connector = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 100, 200, 300, 400)

    assert isinstance(connector, Connector)
    assert connector.shape_id == 2
    assert connector.name == "Connector 1"
    assert snapshot == parent.part._element.cSld.spTree.xml


def test_slide_shapes_add_group_shape(parent, snapshot: SnapshotAssertion) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    group = shapes.add_group_shape()

    assert isinstance(group, GroupShape)
    assert group.shape_id == 2
    assert group.name == "Group 1"
    assert snapshot == parent.part._element.cSld.spTree.xml


def test_slide_shapes_add_table(parent, snapshot: SnapshotAssertion) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    graphic_frame = shapes.add_table(2, 3, 100, 200, 300, 400)

    assert isinstance(graphic_frame, GraphicFrame)
    assert graphic_frame.has_table is True
    assert snapshot == parent.part._element.cSld.spTree.xml


def test_slide_shapes_add_chart(parent, snapshot: SnapshotAssertion) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)
    chart_data = CategoryChartData()
    chart_data.categories = ("A", "B")
    chart_data.add_series("Series 1", (1, 2))

    chart_shape = shapes.add_chart(XL_CHART_TYPE.LINE, 100, 200, 300, 400, chart_data)

    assert chart_shape.has_chart is True
    assert snapshot == parent.part._element.cSld.spTree.xml


def test_slide_shapes_add_picture(
    parent, test_files_dir: Path, snapshot: SnapshotAssertion
) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    picture = shapes.add_picture(str(test_files_dir / "python-icon.jpeg"), 100, 200, 300, 400)

    assert isinstance(picture, Picture)
    assert picture.shape_id == 2
    assert picture.name == "Picture 1"
    assert snapshot == parent.part._element.cSld.spTree.xml


def test_slide_shapes_add_movie(parent, test_files_dir: Path, snapshot: SnapshotAssertion) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    movie = shapes.add_movie(
        str(test_files_dir / "dummy.mp4"),
        100,
        200,
        300,
        400,
        poster_frame_image=str(test_files_dir / "python-icon.jpeg"),
    )

    assert isinstance(movie, Movie)
    assert movie.shape_type is not None
    assert snapshot == parent.part._element.xml


def test_slide_shapes_add_ole_object(parent, snapshot: SnapshotAssertion) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    graphic_frame = shapes.add_ole_object(
        io.BytesIO(b"12345"),
        PROG_ID.XLSX,
        left=100,
        top=200,
        width=None,
        height=None,
    )

    assert isinstance(graphic_frame, GraphicFrame)
    assert graphic_frame.shape_type is not None
    assert snapshot == parent.part._element.cSld.spTree.xml


def test_slide_shapes_title(parent) -> None:
    sp = parse_xml(
        b"""
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="2" name="Title 1"/>
            <p:cNvSpPr/>
            <p:nvPr><p:ph type="title" idx="0"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr><a:prstGeom prst="rect"/></p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
    )
    parent.part._element.cSld.spTree.append(sp)
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    title_shape = shapes.title

    assert title_shape is not None
    assert title_shape.shape_id == 2


def test_slide_shapes_turbo_add(parent) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    assert shapes.turbo_add_enabled is False

    shapes.turbo_add_enabled = True
    first = shapes.add_textbox(10, 10, 10, 10)
    second = shapes.add_textbox(20, 20, 20, 20)

    assert shapes.turbo_add_enabled is True
    assert first.shape_id == 2
    assert second.shape_id == 3


def test_slide_shapes_index(parent) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)
    one = shapes.add_textbox(10, 10, 10, 10)
    two = shapes.add_textbox(20, 20, 20, 20)

    assert shapes.index(one) == 0
    assert shapes.index(two) == 1


def test_base_shape_factory(parent) -> None:
    sp = parse_xml(
        b"""
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr><p:cNvPr id="2" name="Shape"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
          <p:spPr><a:prstGeom prst="rect"/></p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
    )

    shape = BaseShapeFactory(sp, parent)

    assert shape.shape_id == 2


def test_slide_shape_factory_for_placeholder(parent) -> None:
    sp = parse_xml(
        b"""
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="2" name="Title 1"/>
            <p:cNvSpPr/>
            <p:nvPr><p:ph type="title" idx="0"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr><a:prstGeom prst="rect"/></p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
    )

    shape = SlideShapeFactory(sp, parent)

    assert isinstance(shape, SlidePlaceholder)


def test_slide_shape_factory_specialized_placeholders(parent) -> None:
    chart_sp = parse_xml(
        b"""
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="2" name="Chart Placeholder"/>
            <p:cNvSpPr/>
            <p:nvPr><p:ph type="chart" idx="1"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr><a:prstGeom prst="rect"/></p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
    )
    picture_sp = parse_xml(
        b"""
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="3" name="Picture Placeholder"/>
            <p:cNvSpPr/>
            <p:nvPr><p:ph type="pic" idx="2"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr><a:prstGeom prst="rect"/></p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
    )
    table_sp = parse_xml(
        b"""
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="4" name="Table Placeholder"/>
            <p:cNvSpPr/>
            <p:nvPr><p:ph type="tbl" idx="3"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr><a:prstGeom prst="rect"/></p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
    )

    assert isinstance(SlideShapeFactory(chart_sp, parent), ChartPlaceholder)
    assert isinstance(SlideShapeFactory(picture_sp, parent), PicturePlaceholder)
    assert isinstance(SlideShapeFactory(table_sp, parent), TablePlaceholder)


def test_slide_placeholders_collection(parent) -> None:
    sp_tree = parent.part._element.cSld.spTree
    sp_tree.append(
        parse_xml(
            b"""
            <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvSpPr>
                <p:cNvPr id="2" name="Body"/>
                <p:cNvSpPr/>
                <p:nvPr><p:ph type="body" idx="1"/></p:nvPr>
              </p:nvSpPr>
              <p:spPr><a:prstGeom prst="rect"/></p:spPr>
              <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
            </p:sp>
            """
        )
    )
    sp_tree.append(
        parse_xml(
            b"""
            <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvSpPr>
                <p:cNvPr id="3" name="Title"/>
                <p:cNvSpPr/>
                <p:nvPr><p:ph type="title" idx="0"/></p:nvPr>
              </p:nvSpPr>
              <p:spPr><a:prstGeom prst="rect"/></p:spPr>
              <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
            </p:sp>
            """
        )
    )

    placeholders = SlidePlaceholders(sp_tree, parent)

    assert len(placeholders) == 2
    assert placeholders[0].shape_id == 3
    assert placeholders[1].shape_id == 2
    assert [ph.shape_id for ph in placeholders] == [3, 2]


def test_layout_master_notes_shape_factories(
    parent, slide_layout_part, slide_master_part
) -> None:
    sp = parse_xml(
        b"""
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="2" name="Placeholder"/>
            <p:cNvSpPr/>
            <p:nvPr><p:ph type="body" idx="1"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr><a:prstGeom prst="rect"/></p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
    )

    layout_shape = LayoutShapes(
        slide_layout_part._element.cSld.spTree, slide_layout_part.slide_layout
    )._shape_factory(sp)
    master_shape = MasterShapes(
        slide_master_part._element.cSld.spTree, slide_master_part.slide_master
    )._shape_factory(sp)

    notes_xml = parse_xml(
        b"""
        <p:notes xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <p:cSld>
            <p:spTree>
              <p:nvGrpSpPr>
                <p:cNvPr id="1" name=""/>
                <p:cNvGrpSpPr/>
                <p:nvPr/>
              </p:nvGrpSpPr>
              <p:grpSpPr>
                <a:xfrm>
                  <a:off x="0" y="0"/>
                  <a:ext cx="0" cy="0"/>
                  <a:chOff x="0" y="0"/>
                  <a:chExt cx="0" cy="0"/>
                </a:xfrm>
              </p:grpSpPr>
            </p:spTree>
          </p:cSld>
        </p:notes>
        """
    )
    notes_shapes = NotesSlideShapes(notes_xml.cSld.spTree, parent)
    notes_shape = notes_shapes._shape_factory(sp)

    assert isinstance(layout_shape, LayoutPlaceholder)
    assert isinstance(master_shape, MasterPlaceholder)
    assert isinstance(notes_shape, NotesSlidePlaceholder)


def test_layout_master_notes_placeholders_collections(
    parent, slide_layout_part, slide_master_part
) -> None:
    layout_sp = slide_layout_part._element.cSld.spTree
    master_sp = slide_master_part._element.cSld.spTree

    layout_placeholder = parse_xml(
        b"""
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="2" name="Placeholder"/>
            <p:cNvSpPr/>
            <p:nvPr><p:ph type="body" idx="4"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr><a:prstGeom prst="rect"/></p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
    )
    master_placeholder = parse_xml(
        b"""
        <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvSpPr>
            <p:cNvPr id="3" name="Placeholder"/>
            <p:cNvSpPr/>
            <p:nvPr><p:ph type="body" idx="5"/></p:nvPr>
          </p:nvSpPr>
          <p:spPr><a:prstGeom prst="rect"/></p:spPr>
          <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
        </p:sp>
        """
    )
    layout_sp.append(layout_placeholder)
    master_sp.append(master_placeholder)

    layout_placeholders = LayoutPlaceholders(layout_sp, slide_layout_part.slide_layout)
    master_placeholders = MasterPlaceholders(master_sp, slide_master_part.slide_master)
    notes_placeholders = NotesSlidePlaceholders(master_sp, parent)

    assert layout_placeholders.get(idx=4) is not None
    assert master_placeholders.get(PP_PLACEHOLDER.BODY) is not None
    assert len(list(notes_placeholders)) >= 1
