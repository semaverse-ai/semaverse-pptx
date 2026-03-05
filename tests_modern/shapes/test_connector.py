from __future__ import annotations

import pytest
from syrupy.assertion import SnapshotAssertion

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.shapes.autoshape import Shape
from pptx.shapes.connector import Connector


def _connector(
    parent,
    *,
    x: int = 100,
    y: int = 200,
    cx: int = 300,
    cy: int = 400,
    flip_h: bool = False,
    flip_v: bool = False,
) -> Connector:
    flip_h_attr = ' flipH="1"' if flip_h else ""
    flip_v_attr = ' flipV="1"' if flip_v else ""
    xml = """
        <p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:nvCxnSpPr>
            <p:cNvPr id="42" name="Connector 1"/>
            <p:cNvCxnSpPr/>
            <p:nvPr/>
          </p:nvCxnSpPr>
          <p:spPr>
            <a:xfrm{flip_h_attr}{flip_v_attr}>
              <a:off x="{x}" y="{y}"/>
              <a:ext cx="{cx}" cy="{cy}"/>
            </a:xfrm>
          </p:spPr>
        </p:cxnSp>
    """.format(
        flip_h_attr=flip_h_attr,
        flip_v_attr=flip_v_attr,
        x=x,
        y=y,
        cx=cx,
        cy=cy,
    )

    return Connector(
        parse_xml(xml.encode("utf-8")),
        parent,
    )


def _shape_for_connection(parent) -> Shape:
    return Shape(
        parse_xml(
            b"""
            <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvSpPr>
                <p:cNvPr id="77" name="Target"/>
                <p:cNvSpPr/>
                <p:nvPr/>
              </p:nvSpPr>
              <p:spPr>
                <a:xfrm>
                  <a:off x="100" y="200"/>
                  <a:ext cx="300" cy="400"/>
                </a:xfrm>
                <a:prstGeom prst="rect"/>
              </p:spPr>
              <p:txBody>
                <a:bodyPr/>
                <a:lstStyle/>
                <a:p/>
              </p:txBody>
            </p:sp>
            """
        ),
        parent,
    )


def test_connector_point_properties(parent) -> None:
    connector = Connector(
        parse_xml(
            b"""
            <p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                     xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvCxnSpPr>
                <p:cNvPr id="42" name="Connector 1"/>
                <p:cNvCxnSpPr/>
                <p:nvPr/>
              </p:nvCxnSpPr>
              <p:spPr>
                <a:xfrm>
                  <a:off x="100" y="200"/>
                  <a:ext cx="300" cy="400"/>
                </a:xfrm>
              </p:spPr>
            </p:cxnSp>
            """
        ),
        parent,
    )

    assert connector.begin_x == 100
    assert connector.begin_y == 200
    assert connector.end_x == 400
    assert connector.end_y == 600

    connector.begin_x = 150
    connector.begin_y = 250
    connector.end_x = 500
    connector.end_y = 700

    assert connector.begin_x == 150
    assert connector.begin_y == 250
    assert connector.end_x == 500
    assert connector.end_y == 700


def test_connector_connectors_add_st_cxn_and_end_cxn(parent, snapshot: SnapshotAssertion) -> None:
    connector = Connector(
        parse_xml(
            b"""
            <p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                     xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvCxnSpPr>
                <p:cNvPr id="42" name="Connector 1"/>
                <p:cNvCxnSpPr/>
                <p:nvPr/>
              </p:nvCxnSpPr>
              <p:spPr>
                <a:xfrm>
                  <a:off x="10" y="20"/>
                  <a:ext cx="30" cy="40"/>
                </a:xfrm>
              </p:spPr>
            </p:cxnSp>
            """
        ),
        parent,
    )
    shape = _shape_for_connection(parent)

    connector.begin_connect(shape, 0)
    connector.end_connect(shape, 2)

    assert connector.begin_x == 250
    assert connector.begin_y == 200
    assert connector.end_x == 250
    assert connector.end_y == 600
    assert snapshot == connector._element.xml


def test_connector_line_and_shape_type(parent) -> None:
    connector = Connector(
        parse_xml(
            b"""
            <p:cxnSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                     xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvCxnSpPr>
                <p:cNvPr id="42" name="Connector 1"/>
                <p:cNvCxnSpPr/>
                <p:nvPr/>
              </p:nvCxnSpPr>
              <p:spPr>
                <a:xfrm>
                  <a:off x="100" y="200"/>
                  <a:ext cx="300" cy="400"/>
                </a:xfrm>
              </p:spPr>
            </p:cxnSp>
            """
        ),
        parent,
    )

    connector.line.width = 12700

    assert connector.shape_type == MSO_SHAPE_TYPE.LINE
    assert connector.line.width == 12700


@pytest.mark.parametrize(
    ("flip_h", "new_x", "expected_x", "expected_cx", "expected_flip_h"),
    [
        (True, 450, 100, 350, True),
        (True, 150, 100, 50, True),
        (True, 50, 50, 50, False),
        (False, 50, 50, 350, False),
        (False, 200, 200, 200, False),
        (False, 450, 400, 50, True),
    ],
)
def test_connector_begin_x_setter_all_branches(
    parent,
    flip_h: bool,
    new_x: int,
    expected_x: int,
    expected_cx: int,
    expected_flip_h: bool,
) -> None:
    # Arrange
    connector = _connector(parent, flip_h=flip_h)

    # Act
    connector.begin_x = new_x

    # Assert
    assert int(connector._element.x) == expected_x
    assert int(connector._element.cx) == expected_cx
    assert bool(connector._element.flipH) is expected_flip_h


@pytest.mark.parametrize(
    ("flip_v", "new_y", "expected_y", "expected_cy", "expected_flip_v"),
    [
        (True, 650, 200, 450, True),
        (True, 350, 200, 150, True),
        (True, 50, 50, 150, False),
        (False, 50, 50, 550, False),
        (False, 350, 350, 250, False),
        (False, 700, 600, 100, True),
    ],
)
def test_connector_begin_y_setter_all_branches(
    parent,
    flip_v: bool,
    new_y: int,
    expected_y: int,
    expected_cy: int,
    expected_flip_v: bool,
) -> None:
    # Arrange
    connector = _connector(parent, flip_v=flip_v)

    # Act
    connector.begin_y = new_y

    # Assert
    assert int(connector._element.y) == expected_y
    assert int(connector._element.cy) == expected_cy
    assert bool(connector._element.flipV) is expected_flip_v


@pytest.mark.parametrize(
    ("flip_h", "new_x", "expected_x", "expected_cx", "expected_flip_h"),
    [
        (True, 50, 50, 350, True),
        (True, 250, 250, 150, True),
        (True, 500, 400, 100, False),
        (False, 500, 100, 400, False),
        (False, 250, 100, 150, False),
        (False, 50, 50, 50, True),
    ],
)
def test_connector_end_x_setter_all_branches(
    parent,
    flip_h: bool,
    new_x: int,
    expected_x: int,
    expected_cx: int,
    expected_flip_h: bool,
) -> None:
    # Arrange
    connector = _connector(parent, flip_h=flip_h)

    # Act
    connector.end_x = new_x

    # Assert
    assert int(connector._element.x) == expected_x
    assert int(connector._element.cx) == expected_cx
    assert bool(connector._element.flipH) is expected_flip_h


@pytest.mark.parametrize(
    ("flip_v", "new_y", "expected_y", "expected_cy", "expected_flip_v"),
    [
        (True, 150, 150, 450, True),
        (True, 350, 350, 250, True),
        (True, 700, 600, 100, False),
        (False, 700, 200, 500, False),
        (False, 450, 200, 250, False),
        (False, 50, 50, 150, True),
    ],
)
def test_connector_end_y_setter_all_branches(
    parent,
    flip_v: bool,
    new_y: int,
    expected_y: int,
    expected_cy: int,
    expected_flip_v: bool,
) -> None:
    # Arrange
    connector = _connector(parent, flip_v=flip_v)

    # Act
    connector.end_y = new_y

    # Assert
    assert int(connector._element.y) == expected_y
    assert int(connector._element.cy) == expected_cy
    assert bool(connector._element.flipV) is expected_flip_v
