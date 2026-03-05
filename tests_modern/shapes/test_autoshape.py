from __future__ import annotations

import pytest

from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.shapes.autoshape import Adjustment, AdjustmentCollection, AutoShapeType, Shape


@pytest.mark.parametrize(
    ("def_val", "actual", "expected"),
    [
        (50000, None, 0.5),
        (50000, 12500, 0.125),
        (0, None, 0.0),
        (-20833, None, -0.20833),
    ],
)
def test_adjustment_effective_value(def_val: int, actual: int | None, expected: float) -> None:
    assert Adjustment("adj", def_val, actual).effective_value == expected


def test_adjustment_effective_value_setter() -> None:
    adj = Adjustment("adj", 50000)

    adj.effective_value = 0.25

    assert adj.actual == 25000
    assert adj.val == 25000


@pytest.mark.parametrize("bad_value", ["x", object()])
def test_adjustment_effective_value_setter_raises_on_non_numeric(bad_value: object) -> None:
    adj = Adjustment("adj", 50000)

    with pytest.raises(ValueError, match="adjustment value must be numeric"):
        adj.effective_value = bad_value  # type: ignore[assignment]


def test_adjustment_collection_indexing_roundtrip() -> None:
    prst_geom = parse_xml(
        b"""
        <a:prstGeom xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="chevron">
          <a:avLst>
            <a:gd name="adj" fmla="val 25000"/>
          </a:avLst>
        </a:prstGeom>
        """
    )

    adjustments = AdjustmentCollection(prst_geom)

    assert len(adjustments) == 1
    assert adjustments[0] == 0.25

    adjustments[0] = 0.5

    assert adjustments[0] == 0.5
    assert prst_geom.avLst.gd_lst[0].fmla == "val 50000"


def test_autoshape_type_properties() -> None:
    auto = AutoShapeType(MSO_SHAPE.ROUNDED_RECTANGLE)

    assert auto.autoshape_type_id == MSO_SHAPE.ROUNDED_RECTANGLE
    assert auto.basename == "Rounded Rectangle"
    assert auto.prst == "roundRect"


def test_autoshape_type_special_character_basename() -> None:
    auto = AutoShapeType(MSO_SHAPE.NO_SYMBOL)

    assert auto.basename == "&quot;No&quot; Symbol"


def test_autoshape_type_mapping_helpers() -> None:
    assert AutoShapeType.id_from_prst("rect") == MSO_SHAPE.RECTANGLE
    assert AutoShapeType.default_adjustment_values(MSO_SHAPE.CHEVRON) == (("adj", 50000),)


def test_autoshape_type_raises_on_unknown_type() -> None:
    with pytest.raises(KeyError, match="no autoshape type"):
        AutoShapeType(99_999)  # type: ignore[arg-type]


def test_shape_properties_and_text(parent) -> None:
    shape = Shape(
        parse_xml(
            b"""
            <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvSpPr>
                <p:cNvPr id="42" name="Shape 1"/>
                <p:cNvSpPr/>
                <p:nvPr/>
              </p:nvSpPr>
              <p:spPr>
                <a:prstGeom prst="rect"/>
              </p:spPr>
              <p:txBody>
                <a:bodyPr/>
                <a:lstStyle/>
                <a:p>
                  <a:r><a:t>hello</a:t></a:r>
                </a:p>
              </p:txBody>
            </p:sp>
            """
        ),
        parent,
    )

    assert shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert shape.auto_shape_type == MSO_SHAPE.RECTANGLE
    assert shape.has_text_frame is True
    assert shape.text == "hello"

    shape.text = "updated"

    assert shape.text == "updated"


def test_shape_adjustments(parent) -> None:
    shape = Shape(
        parse_xml(
            b"""
            <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvSpPr>
                <p:cNvPr id="42" name="Chevron 1"/>
                <p:cNvSpPr/>
                <p:nvPr/>
              </p:nvSpPr>
              <p:spPr>
                <a:prstGeom prst="chevron">
                  <a:avLst>
                    <a:gd name="adj" fmla="val 25000"/>
                  </a:avLst>
                </a:prstGeom>
              </p:spPr>
              <p:txBody>
                <a:bodyPr/>
                <a:lstStyle/>
                <a:p/>
              </p:txBody>
            </p:sp>
            """
        ),
        parent,
    )

    assert shape.adjustments[0] == 0.25

    shape.adjustments[0] = 0.35

    assert shape.adjustments[0] == 0.35
    assert "val 35000" in shape._element.xml


def test_shape_fill_and_line_access(parent) -> None:
    shape = Shape(
        parse_xml(
            b"""
            <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:nvSpPr>
                <p:cNvPr id="42" name="Shape 1"/>
                <p:cNvSpPr/>
                <p:nvPr/>
              </p:nvSpPr>
              <p:spPr>
                <a:prstGeom prst="rect"/>
              </p:spPr>
              <p:txBody>
                <a:bodyPr/>
                <a:lstStyle/>
                <a:p/>
              </p:txBody>
            </p:sp>
            """
        ),
        parent,
    )

    shape.line.width = 12700

    assert shape.fill is not None
    assert shape.line.width == 12700
