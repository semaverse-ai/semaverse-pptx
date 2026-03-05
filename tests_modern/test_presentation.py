from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO

from syrupy.assertion import SnapshotAssertion

from pptx.oxml import parse_xml
from pptx.presentation import Presentation
from tests_modern.factories import presentation_xml
from tests_modern.stubs import PresentationPartStub
from tests_modern.xml_utils import serialize_xml


def test_presentation_delegates_core_properties_notes_master_and_save() -> None:
    # Arrange
    part = PresentationPartStub()
    prs = Presentation(parse_xml(presentation_xml()), part)  # type: ignore[arg-type]

    # Act
    stream = BytesIO()
    prs.save(stream)

    # Assert
    assert prs.core_properties == "core-props"
    assert prs.notes_master == "notes-master"
    assert stream.getvalue() == b"saved"
    assert part.saved_payload == b"saved"


def test_presentation_slide_width_and_height_getters_and_setters() -> None:
    # Arrange
    xml = presentation_xml(b'<p:sldSz cx="9144000" cy="6858000"/>')
    prs = Presentation(parse_xml(xml), PresentationPartStub())  # type: ignore[arg-type]

    # Act
    prs.slide_width = 1000000
    prs.slide_height = 2000000

    # Assert
    assert prs.slide_width == 1000000
    assert prs.slide_height == 2000000


def test_presentation_slide_size_setters_create_sldsz_when_missing(
    snapshot: SnapshotAssertion,
) -> None:
    # Arrange
    prs = Presentation(parse_xml(presentation_xml()), PresentationPartStub())  # type: ignore[arg-type]

    # Act
    prs.slide_width = 914400
    prs.slide_height = 1828800

    # Assert
    assert serialize_xml(prs._element) == snapshot


def test_presentation_slides_calls_rename_slide_parts_with_rids() -> None:
    # Arrange
    xml = presentation_xml(
        b"""
        <p:sldIdLst>
          <p:sldId id="256" r:id="rId1"/>
          <p:sldId id="257" r:id="rId2"/>
        </p:sldIdLst>
        """
    )
    part = PresentationPartStub(
        slides_by_rid={
            "rId1": "slide-1",
            "rId2": "slide-2",
        }
    )
    prs = Presentation(parse_xml(xml), part)  # type: ignore[arg-type]

    # Act
    slides = prs.slides

    # Assert
    assert part.renamed_rids == ["rId1", "rId2"]
    assert len(slides) == 2
    assert slides[0] == "slide-1"
    assert list(slides) == ["slide-1", "slide-2"]


def test_presentation_slides_creates_sldidlst_when_missing(snapshot: SnapshotAssertion) -> None:
    # Arrange
    part = PresentationPartStub()
    prs = Presentation(parse_xml(presentation_xml()), part)  # type: ignore[arg-type]

    # Act
    slides = prs.slides

    # Assert
    assert part.renamed_rids == []
    assert len(slides) == 0
    assert serialize_xml(prs._element) == snapshot


def test_presentation_slide_masters_and_slide_master() -> None:
    # Arrange
    xml = presentation_xml(
        b"""
        <p:sldMasterIdLst>
          <p:sldMasterId r:id="rId1"/>
          <p:sldMasterId r:id="rId2"/>
        </p:sldMasterIdLst>
        """
    )
    master_1 = object()
    master_2 = object()
    part = PresentationPartStub(slide_masters_by_rid={"rId1": master_1, "rId2": master_2})
    prs = Presentation(parse_xml(xml), part)  # type: ignore[arg-type]

    # Act
    masters = prs.slide_masters

    # Assert
    assert len(masters) == 2
    assert masters[0] is master_1
    assert list(masters) == [master_1, master_2]
    assert prs.slide_master is master_1


def test_presentation_slide_layouts_comes_from_first_master() -> None:
    # Arrange
    xml = presentation_xml(
        b"""
        <p:sldMasterIdLst>
          <p:sldMasterId r:id="rId1"/>
        </p:sldMasterIdLst>
        """
    )

    @dataclass
    class _SlideMasterStub:
        slide_layouts: object

    expected_layouts = ("layout-1", "layout-2")
    part = PresentationPartStub(
        slide_masters_by_rid={"rId1": _SlideMasterStub(slide_layouts=expected_layouts)}
    )
    prs = Presentation(parse_xml(xml), part)  # type: ignore[arg-type]

    # Act
    layouts = prs.slide_layouts

    # Assert
    assert layouts is expected_layouts


def test_presentation_slide_masters_creates_list_when_missing(
    snapshot: SnapshotAssertion,
) -> None:
    # Arrange
    prs = Presentation(parse_xml(presentation_xml()), PresentationPartStub())  # type: ignore[arg-type]

    # Act
    masters = prs.slide_masters

    # Assert
    assert len(masters) == 0
    assert serialize_xml(prs._element) == snapshot
