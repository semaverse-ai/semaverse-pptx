from __future__ import annotations

# ruff: noqa: E501
from typing import TYPE_CHECKING

from pptx.chart.chart import Chart, ChartTitle, _Plots
from pptx.chart.data import CategoryChartData
from pptx.chart.legend import Legend
from pptx.chart.series import SeriesCollection
from pptx.dml.chtfmt import ChartFormat
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml import parse_xml

if TYPE_CHECKING:
    from tests_modern.chart.conftest import ChartPartStub


def _chart_space(xml_body: bytes = b"") -> object:
    return parse_xml(
        b'<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
        b'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        + xml_body
        + b"</c:chartSpace>"
    )


def test_chart_core_properties(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(
            b"<c:chart><c:title><c:layout/><c:tx><c:rich/></c:tx></c:title>"
            b"<c:plotArea><c:pieChart/><c:catAx/><c:valAx/></c:plotArea><c:legend/></c:chart>"
        ),
        chart_part_stub,
    )

    assert chart.has_title is True
    assert isinstance(chart.chart_title, ChartTitle)
    assert chart.has_legend is True
    assert isinstance(chart.legend, Legend)
    assert isinstance(chart.plots, _Plots)
    assert isinstance(chart.series, SeriesCollection)
    assert chart.chart_type == XL_CHART_TYPE.PIE


def test_chart_style_and_legend_toggle(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(
            b"<c:chart><c:plotArea><c:barChart/><c:catAx/><c:valAx/></c:plotArea><c:legend/></c:chart>"
        ),
        chart_part_stub,
    )

    assert chart.chart_style is None

    chart.chart_style = 7

    assert chart.chart_style == 7

    chart.has_legend = False

    assert chart.has_legend is False


def test_chart_replace_data(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(
            b"<c:chart><c:plotArea><c:barChart><c:barDir val='col'/><c:grouping val='clustered'/>"
            b"<c:ser><c:idx val='0'/><c:order val='0'/></c:ser></c:barChart>"
            b"<c:catAx><c:axId val='1'/><c:crossAx val='2'/></c:catAx>"
            b"<c:valAx><c:axId val='2'/><c:crossAx val='1'/></c:valAx>"
            b"</c:plotArea></c:chart>"
        ),
        chart_part_stub,
    )
    data = CategoryChartData()
    data.categories = ("A", "B")
    data.add_series("Series 1", (1, 2))

    chart.replace_data(data)

    assert chart_part_stub.chart_workbook.updated_blob is not None


def test_chart_title_behaviors() -> None:
    title = ChartTitle(
        parse_xml(b'<c:title xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"/>')
    )

    assert title.has_text_frame is False

    title.has_text_frame = True

    assert title.has_text_frame is True
    assert title.text_frame is not None
    assert isinstance(title.format, ChartFormat)


def test_legend_behaviors() -> None:
    legend = Legend(
        parse_xml(b'<c:legend xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"/>')
    )

    assert legend.include_in_layout is True
    assert legend.position == XL_LEGEND_POSITION.RIGHT

    legend.horz_offset = 0.5
    legend.include_in_layout = False
    legend.position = XL_LEGEND_POSITION.BOTTOM

    assert legend.horz_offset == 0.5
    assert legend.include_in_layout is False
    assert legend.position == XL_LEGEND_POSITION.BOTTOM
