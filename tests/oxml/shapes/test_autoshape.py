from __future__ import annotations

import pytest

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml import parse_xml
from pptx.oxml.shapes.autoshape import CT_Shape
from pptx.oxml.shapes.shared import ST_Direction, ST_PlaceholderSize


def test_prst_geom_gd_lst() -> None:
    prst_geom = parse_xml(
        '<a:prstGeom xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'prst="chevron"><a:avLst><a:gd name="adj1" fmla="val 111"/>'
        '<a:gd name="adj2" fmla="val 222"/></a:avLst></a:prstGeom>'
    )

    gd_vals = [(gd.name, gd.fmla) for gd in prst_geom.gd_lst]

    assert gd_vals == [("adj1", "val 111"), ("adj2", "val 222")]


def test_prst_geom_rewrite_guides() -> None:
    prst_geom = parse_xml(
        '<a:prstGeom xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'prst="chevron"><a:avLst><a:gd name="adj6" fmla="val 666"/></a:avLst></a:prstGeom>'
    )

    prst_geom.rewrite_guides([("adj1", 111), ("adj2", 222)])

    assert [gd.fmla for gd in prst_geom.gd_lst] == ["val 111", "val 222"]


def test_shape_new_autoshape_sp() -> None:
    sp = CT_Shape.new_autoshape_sp(9, "Rounded Rectangle 8", "roundRect", 111, 222, 333, 444)

    assert sp.nvSpPr.cNvPr.id == 9
    assert sp.nvSpPr.cNvPr.name == "Rounded Rectangle 8"
    assert sp.prst == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE


@pytest.mark.parametrize(
    ("id_", "name", "ph_type", "orient", "sz", "idx"),
    [
        (
            2,
            "Title 1",
            PP_PLACEHOLDER.CENTER_TITLE,
            ST_Direction.HORZ,
            ST_PlaceholderSize.FULL,
            0,
        ),
        (
            4,
            "Vertical Subtitle 3",
            PP_PLACEHOLDER.SUBTITLE,
            ST_Direction.VERT,
            ST_PlaceholderSize.FULL,
            1,
        ),
        (
            7,
            "Footer Placeholder 6",
            PP_PLACEHOLDER.FOOTER,
            ST_Direction.HORZ,
            ST_PlaceholderSize.QUARTER,
            11,
        ),
    ],
)
def test_shape_new_placeholder_sp(
    id_: int,
    name: str,
    ph_type: PP_PLACEHOLDER,
    orient: ST_Direction | str,
    sz: ST_PlaceholderSize | str,
    idx: int,
) -> None:
    sp = CT_Shape.new_placeholder_sp(id_, name, ph_type, orient, sz, idx)

    assert sp.nvSpPr.cNvPr.id == id_
    assert sp.nvSpPr.cNvPr.name == name
    assert sp.nvSpPr.nvPr.ph.type == ph_type


def test_shape_new_textbox_sp() -> None:
    sp = CT_Shape.new_textbox_sp(9, "TextBox 8", 111, 222, 333, 444)

    assert sp.nvSpPr.cNvPr.id == 9
    assert sp.nvSpPr.cNvPr.name == "TextBox 8"
    assert sp.nvSpPr.cNvSpPr.txBox is True


@pytest.mark.parametrize(
    ("xml", "expected"),
    [
        (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            "<p:nvSpPr><p:cNvSpPr/></p:nvSpPr><p:spPr><a:prstGeom/></p:spPr></p:sp>",
            True,
        ),
        (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            "<p:nvSpPr><p:nvPr><p:ph/></p:nvPr></p:nvSpPr><p:spPr/></p:sp>",
            False,
        ),
        (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<p:nvSpPr><p:cNvSpPr txBox="1"/></p:nvSpPr><p:spPr><a:prstGeom/></p:spPr></p:sp>',
            False,
        ),
    ],
)
def test_shape_is_autoshape(xml: str, expected: bool) -> None:
    sp = parse_xml(xml)

    assert sp.is_autoshape is expected


@pytest.mark.parametrize(
    ("xml", "expected"),
    [
        (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            "<p:nvSpPr><p:cNvSpPr/></p:nvSpPr><p:spPr><a:prstGeom/></p:spPr></p:sp>",
            False,
        ),
        (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            "<p:nvSpPr><p:cNvSpPr/><p:nvPr><p:ph/></p:nvPr></p:nvSpPr><p:spPr/></p:sp>",
            False,
        ),
        (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<p:nvSpPr><p:cNvSpPr txBox="1"/></p:nvSpPr><p:spPr><a:prstGeom/></p:spPr></p:sp>',
            True,
        ),
    ],
)
def test_shape_is_textbox(xml: str, expected: bool) -> None:
    sp = parse_xml(xml)

    assert sp.is_textbox is expected


def test_shape_add_path() -> None:
    sp = parse_xml(
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        "<p:spPr><a:custGeom/></p:spPr></p:sp>"
    )

    path = sp.add_path(100, 200)

    assert path.tag.endswith("path")
    assert len(sp.spPr.custGeom.pathLst) == 1
    assert sp.spPr.custGeom.pathLst[0].w == 100


def test_shape_add_path_raises_if_not_freeform() -> None:
    sp = parse_xml(
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><p:spPr/></p:sp>'
    )

    with pytest.raises(ValueError, match="shape must be freeform"):
        sp.add_path(100, 200)


def test_shape_get_or_add_ln() -> None:
    sp = parse_xml(
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><p:spPr/></p:sp>'
    )

    assert sp.ln is None

    ln = sp.get_or_add_ln()

    assert sp.ln is ln
    assert sp.ln is not None


def test_shape_prst_value() -> None:
    sp = parse_xml(
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:spPr><a:prstGeom prst="rect"/></p:spPr></p:sp>'
    )

    assert sp.prst == MSO_AUTO_SHAPE_TYPE.RECTANGLE


def test_shape_prst_none() -> None:
    sp = parse_xml(
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><p:spPr/></p:sp>'
    )

    assert sp.prst is None


def test_path2d_operations() -> None:
    path = parse_xml('<a:path xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')

    path.add_moveTo(10, 20)
    path.add_lnTo(30, 40)
    path.add_close()

    assert len(path) == 3
    assert path[0].tag.endswith("moveTo")
    assert path[1].tag.endswith("lnTo")
    assert path[2].tag.endswith("close")
