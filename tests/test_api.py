from __future__ import annotations

from pathlib import Path

import pytest

from pptx.api import Presentation, _default_pptx_path
from pptx.exc import PackageNotFoundError
from pptx.opc.constants import CONTENT_TYPE as CT
from tests.stubs import PackageStub, PresentationPartStub


def test_presentation_uses_default_path_when_no_argument(monkeypatch: pytest.MonkeyPatch) -> None:
    # Arrange
    default_path = "/tmp/default.pptx"
    expected_prs = object()
    seen: dict[str, str] = {}

    def fake_open(path: str) -> PackageStub:
        seen["path"] = path
        return PackageStub(PresentationPartStub(CT.PML_PRESENTATION_MAIN, expected_prs))

    monkeypatch.setattr("pptx.api._default_pptx_path", lambda: default_path)
    monkeypatch.setattr("pptx.api.Package.open", fake_open)

    # Act
    prs = Presentation()

    # Assert
    assert seen["path"] == default_path
    assert prs is expected_prs


@pytest.mark.parametrize(
    "content_type",
    [CT.PML_PRESENTATION_MAIN, CT.PML_PRES_MACRO_MAIN],
)
def test_presentation_accepts_valid_main_document_content_type(
    monkeypatch: pytest.MonkeyPatch, content_type: str
) -> None:
    # Arrange
    expected_prs = object()
    monkeypatch.setattr(
        "pptx.api.Package.open",
        lambda _: PackageStub(PresentationPartStub(content_type, expected_prs)),
    )

    # Act
    prs = Presentation("dummy.pptx")

    # Assert
    assert prs is expected_prs


def test_presentation_rejects_invalid_main_document_content_type(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    # Arrange
    monkeypatch.setattr(
        "pptx.api.Package.open",
        lambda _: PackageStub(PresentationPartStub("application/not-pptx", object())),
    )

    # Act / Assert
    with pytest.raises(ValueError, match="is not a PowerPoint file"):
        Presentation("dummy.pptx")


def test_default_pptx_path_points_to_builtin_template() -> None:
    # Act
    path = _default_pptx_path()

    # Assert
    assert Path(path).is_absolute()
    assert path.endswith("templates/default.pptx")


def test_presentation_raises_on_missing_path(tmp_path: Path) -> None:
    # Arrange
    missing = tmp_path / "does-not-exist.pptx"

    # Act / Assert
    with pytest.raises(PackageNotFoundError, match="Package not found"):
        Presentation(str(missing))
