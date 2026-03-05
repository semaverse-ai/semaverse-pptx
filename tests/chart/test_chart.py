from __future__ import annotations

# ruff: noqa: E501
from typing import TYPE_CHECKING

import pytest

from pptx.chart.axis import CategoryAxis, DateAxis, ValueAxis
from pptx.chart.chart import Chart, ChartTitle, _Plots
from pptx.chart.data import CategoryChartData
from pptx.chart.legend import Legend
from pptx.chart.series import SeriesCollection
from pptx.dml.chtfmt import ChartFormat
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml import parse_xml
from pptx.text.text import Font
from tests.stubs import PartProviderStub

if TYPE_CHECKING:
    from tests.chart.conftest import ChartPartStub


def _chart_space(xml_body: bytes = b"") -> object:
    return parse_xml(
        b'<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
        b'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        + xml_body
        + b"</c:chartSpace>"
    )


def test_chart_core_properties(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(
            b"<c:chart><c:title><c:layout/><c:tx><c:rich/></c:tx></c:title>"
            b"<c:plotArea><c:pieChart/><c:catAx/><c:valAx/></c:plotArea><c:legend/></c:chart>"
        ),
        chart_part_stub,
    )

    assert chart.has_title is True
    assert isinstance(chart.chart_title, ChartTitle)
    assert chart.has_legend is True
    assert isinstance(chart.legend, Legend)
    assert isinstance(chart.plots, _Plots)
    assert isinstance(chart.series, SeriesCollection)
    assert chart.chart_type == XL_CHART_TYPE.PIE


@pytest.mark.parametrize(
    ("plot_area_xml", "axis_type"),
    [
        (b"<c:plotArea><c:barChart/><c:catAx/></c:plotArea>", CategoryAxis),
        (b"<c:plotArea><c:lineChart/><c:dateAx/></c:plotArea>", DateAxis),
        (b"<c:plotArea><c:bubbleChart/><c:valAx/></c:plotArea>", ValueAxis),
    ],
)
def test_chart_category_axis_fallback_order(
    chart_part_stub: ChartPartStub, plot_area_xml: bytes, axis_type: type[object]
) -> None:
    chart = Chart(_chart_space(b"<c:chart>" + plot_area_xml + b"</c:chart>"), chart_part_stub)

    axis = chart.category_axis

    assert isinstance(axis, axis_type)


def test_chart_category_axis_raises_when_missing(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(b"<c:chart><c:plotArea><c:pieChart/></c:plotArea></c:chart>"), chart_part_stub
    )

    # Act / Assert
    with pytest.raises(ValueError, match="chart has no category axis"):
        _ = chart.category_axis


def test_chart_style_and_legend_toggle(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(
            b"<c:chart><c:plotArea><c:barChart/><c:catAx/><c:valAx/></c:plotArea><c:legend/></c:chart>"
        ),
        chart_part_stub,
    )

    assert chart.chart_style is None

    chart.chart_style = 7

    assert chart.chart_style == 7

    chart.has_legend = False

    assert chart.has_legend is False


def test_chart_style_can_be_cleared(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(
            b"<c:chart><c:plotArea><c:barChart/><c:catAx/><c:valAx/></c:plotArea></c:chart>"
        ),
        chart_part_stub,
    )
    chart.chart_style = 6

    chart.chart_style = None

    assert chart.chart_style is None


def test_chart_replace_data(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(
            b"<c:chart><c:plotArea><c:barChart><c:barDir val='col'/><c:grouping val='clustered'/>"
            b"<c:ser><c:idx val='0'/><c:order val='0'/></c:ser></c:barChart>"
            b"<c:catAx><c:axId val='1'/><c:crossAx val='2'/></c:catAx>"
            b"<c:valAx><c:axId val='2'/><c:crossAx val='1'/></c:valAx>"
            b"</c:plotArea></c:chart>"
        ),
        chart_part_stub,
    )
    data = CategoryChartData()
    data.categories = ("A", "B")
    data.add_series("Series 1", (1, 2))

    chart.replace_data(data)

    assert chart_part_stub.chart_workbook.updated_blob is not None


def test_chart_font_returns_font_instance(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(b"<c:chart><c:plotArea><c:pieChart/></c:plotArea></c:chart>"), chart_part_stub
    )

    font = chart.font

    assert isinstance(font, Font)
    assert chart.font is font


def test_chart_title_and_legend_absence_paths(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(
            b"<c:chart><c:plotArea><c:barChart/><c:catAx/><c:valAx/></c:plotArea></c:chart>"
        ),
        chart_part_stub,
    )

    # Act / Assert
    assert chart.has_title is False
    assert chart.legend is None


def test_chart_has_title_false_removes_title(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(
            b"<c:chart><c:title><c:tx><c:rich/></c:tx></c:title><c:plotArea><c:barChart/><c:catAx/><c:valAx/></c:plotArea></c:chart>"
        ),
        chart_part_stub,
    )

    chart.has_title = False

    assert chart.has_title is False
    assert chart._chartSpace.chart.autoTitleDeleted.val is True


def test_chart_value_axis_paths(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(
            b"""
            <c:chart>
              <c:plotArea>
                <c:barChart/>
                <c:valAx><c:axId val="1"/></c:valAx>
                <c:valAx><c:axId val="2"/></c:valAx>
              </c:plotArea>
            </c:chart>
            """
        ),
        chart_part_stub,
    )
    no_value_axis = Chart(
        _chart_space(b"<c:chart><c:plotArea><c:pieChart/></c:plotArea></c:chart>"),
        chart_part_stub,
    )

    # Act / Assert
    assert 'val="2"' in chart.value_axis._element.xml
    with pytest.raises(ValueError, match="chart has no value axis"):
        _ = no_value_axis.value_axis


def test_chart_title_behaviors() -> None:
    title = ChartTitle(
        parse_xml(b'<c:title xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"/>')
    )

    assert title.has_text_frame is False

    title.has_text_frame = True

    assert title.has_text_frame is True
    assert title.text_frame is not None
    assert isinstance(title.format, ChartFormat)

    title.has_text_frame = False

    assert title.has_text_frame is False


def test_chart_title_exposes_part_from_underlying_element() -> None:
    part = object()
    title = ChartTitle(PartProviderStub(part=part))

    resolved_part = title.part

    assert resolved_part is part


def test_legend_behaviors() -> None:
    legend = Legend(
        parse_xml(b'<c:legend xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"/>')
    )

    assert legend.include_in_layout is True
    assert legend.position == XL_LEGEND_POSITION.RIGHT

    legend.horz_offset = 0.5
    legend.include_in_layout = False
    legend.position = XL_LEGEND_POSITION.BOTTOM

    assert legend.horz_offset == 0.5
    assert legend.include_in_layout is False
    assert legend.position == XL_LEGEND_POSITION.BOTTOM


def test_plots_slice_and_len(chart_part_stub: ChartPartStub) -> None:
    chart = Chart(
        _chart_space(b"<c:chart><c:plotArea><c:barChart/><c:lineChart/></c:plotArea></c:chart>"),
        chart_part_stub,
    )

    plots = chart.plots
    first_plot_slice = plots[:1]

    assert len(plots) == 2
    assert isinstance(first_plot_slice, list)
    assert len(first_plot_slice) == 1
