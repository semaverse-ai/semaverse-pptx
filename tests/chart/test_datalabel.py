from __future__ import annotations

from pptx.chart.datalabel import DataLabel, DataLabels
from pptx.enum.chart import XL_DATA_LABEL_POSITION
from pptx.oxml import parse_xml
from tests.stubs import PartProviderStub


def test_data_labels_properties() -> None:
    labels = DataLabels(
        parse_xml(b'<c:dLbls xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"/>')
    )

    assert labels.number_format == "General"
    assert labels.number_format_is_linked is True
    assert labels.position is None

    labels.number_format = "0.00"
    labels.number_format_is_linked = False
    labels.position = XL_DATA_LABEL_POSITION.INSIDE_END
    labels.show_category_name = True
    labels.show_legend_key = True
    labels.show_percentage = True
    labels.show_series_name = True
    labels.show_value = True

    assert labels.number_format == "0.00"
    assert labels.number_format_is_linked is False
    assert labels.position == XL_DATA_LABEL_POSITION.INSIDE_END
    assert labels.show_category_name is True
    assert labels.show_legend_key is True
    assert labels.show_percentage is True
    assert labels.show_series_name is True
    assert labels.show_value is True


def test_data_label_properties() -> None:
    ser = parse_xml(
        b'<c:ser xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">'
        b"<c:idx val='0'/><c:order val='0'/></c:ser>"
    )
    label = DataLabel(ser, 0)

    assert label.has_text_frame is False
    assert label.position is None

    label.has_text_frame = True
    label.position = XL_DATA_LABEL_POSITION.CENTER

    assert label.has_text_frame is True
    assert label.position == XL_DATA_LABEL_POSITION.CENTER
    assert label.text_frame is not None
    assert label.font is not None

    label.has_text_frame = False

    assert label.has_text_frame is False


def test_data_label_exposes_part_from_series() -> None:
    part = object()
    label = DataLabel(PartProviderStub(part=part), 0)

    resolved_part = label.part

    assert resolved_part is part
