from __future__ import annotations

from typing import TYPE_CHECKING

import pytest
from syrupy.assertion import SnapshotAssertion

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, MSO_UNDERLINE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.shapes.autoshape import Shape
from pptx.text.text import Font, TextFrame, _Paragraph, _Run
from pptx.util import Inches, Pt
from tests.text.factories import paragraph, run, tx_body
from tests.xml_utils import serialize_xml

if TYPE_CHECKING:
    from tests.text.conftest import DummyParent, FitParent


def test_text_frame_text_round_trip() -> None:
    text_frame = TextFrame(
        tx_body(
            b"<a:bodyPr/><a:p><a:r><a:t>Hello</a:t></a:r></a:p><a:p><a:r><a:t>World</a:t></a:r></a:p>"
        ),
        None,
    )

    assert text_frame.text == "Hello\nWorld"

    text_frame.text = "New\nText\vHere"

    assert len(text_frame.paragraphs) == 2
    assert text_frame.paragraphs[0].text == "New"
    assert text_frame.paragraphs[1].text == "Text\vHere"


def test_text_frame_add_paragraph(snapshot: SnapshotAssertion) -> None:
    text_frame = TextFrame(tx_body(b"<a:bodyPr/><a:p/>"), None)

    paragraph_obj = text_frame.add_paragraph()

    assert isinstance(paragraph_obj, _Paragraph)
    assert serialize_xml(text_frame._element) == snapshot


def test_text_frame_clear() -> None:
    text_frame = TextFrame(
        tx_body(b"<a:bodyPr/><a:p><a:r><a:t>One</a:t></a:r></a:p><a:p><a:r><a:t>Two</a:t></a:r></a:p>"),
        None,
    )

    text_frame.clear()

    assert len(text_frame.paragraphs) == 1
    assert text_frame.text == ""


@pytest.mark.parametrize(
    ("xml_body", "expected_value"),
    [
        (b"<a:bodyPr/><a:p/>", None),
        (b"<a:bodyPr><a:noAutofit/></a:bodyPr><a:p/>", MSO_AUTO_SIZE.NONE),
        (b"<a:bodyPr><a:spAutoFit/></a:bodyPr><a:p/>", MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT),
        (b"<a:bodyPr><a:normAutofit/></a:bodyPr><a:p/>", MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE),
    ],
)
def test_text_frame_auto_size_get(xml_body: bytes, expected_value: MSO_AUTO_SIZE | None) -> None:
    text_frame = TextFrame(tx_body(xml_body), None)

    assert text_frame.auto_size == expected_value


def test_text_frame_auto_size_set(snapshot: SnapshotAssertion) -> None:
    text_frame = TextFrame(tx_body(b"<a:bodyPr/><a:p/>"), None)

    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    assert serialize_xml(text_frame._element) == snapshot(name="none")

    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    assert serialize_xml(text_frame._element) == snapshot(name="shape")

    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    assert serialize_xml(text_frame._element) == snapshot(name="text")

    text_frame.auto_size = None
    assert serialize_xml(text_frame._element) == snapshot(name="clear")


def test_text_frame_margins() -> None:
    text_frame = TextFrame(tx_body(b"<a:bodyPr/><a:p/>"), None)

    assert text_frame.margin_left == 91440
    assert text_frame.margin_right == 91440
    assert text_frame.margin_top == 45720
    assert text_frame.margin_bottom == 45720

    text_frame.margin_left = Inches(1)
    text_frame.margin_right = Inches(2)
    text_frame.margin_top = Inches(3)
    text_frame.margin_bottom = Inches(4)

    assert text_frame.margin_left == 914400
    assert text_frame.margin_right == 1828800
    assert text_frame.margin_top == 2743200
    assert text_frame.margin_bottom == 3657600


def test_text_frame_margin_raises_on_invalid_type() -> None:
    text_frame = TextFrame(tx_body(b"<a:bodyPr/><a:p/>"), None)

    with pytest.raises(TypeError):
        setattr(text_frame, "margin_bottom", "100")


def test_text_frame_vertical_anchor() -> None:
    text_frame = TextFrame(tx_body(b"<a:bodyPr/><a:p/>"), None)

    assert text_frame.vertical_anchor is None

    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    assert text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE

    text_frame.vertical_anchor = None
    assert text_frame.vertical_anchor is None


def test_text_frame_word_wrap() -> None:
    text_frame = TextFrame(tx_body(b"<a:bodyPr/><a:p/>"), None)

    assert text_frame.word_wrap is None

    text_frame.word_wrap = True
    assert text_frame.word_wrap is True

    text_frame.word_wrap = False
    assert text_frame.word_wrap is False

    text_frame.word_wrap = None
    assert text_frame.word_wrap is None


def test_text_frame_word_wrap_raises_on_invalid_value() -> None:
    text_frame = TextFrame(tx_body(b"<a:bodyPr/><a:p/>"), None)

    with pytest.raises(ValueError):
        text_frame.word_wrap = "invalid"  # type: ignore[assignment]


def test_text_frame_fit_text_noop_when_empty(monkeypatch: pytest.MonkeyPatch) -> None:
    text_frame = TextFrame(tx_body(b"<a:bodyPr/><a:p/>"), None)
    called = {"flag": False}

    def _unexpected(*_: object, **__: object) -> int:
        called["flag"] = True
        return 10

    monkeypatch.setattr(TextFrame, "_best_fit_font_size", _unexpected)

    text_frame.fit_text()

    assert called["flag"] is False


def test_text_frame_fit_text_applies_best_fit(
    monkeypatch: pytest.MonkeyPatch, fit_parent: FitParent
) -> None:
    text_frame = TextFrame(
        tx_body(b"<a:bodyPr/><a:p><a:r><a:t>Hello world</a:t></a:r></a:p>"),
        fit_parent,
    )

    monkeypatch.setattr("pptx.text.text.FontFiles.find", lambda *args: "font.ttf")
    monkeypatch.setattr("pptx.text.text.TextFitter.best_fit_font_size", lambda *args: 14)

    text_frame.fit_text("Calibri", 40, bold=True, italic=False)

    assert text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert text_frame.word_wrap is True
    assert text_frame.paragraphs[0].runs[0].font.name == "Calibri"
    assert text_frame.paragraphs[0].runs[0].font.size == Pt(14)


def test_text_frame_effective_extents() -> None:
    shape = Shape(
        parse_xml(
            b"<p:sp xmlns:p='http://schemas.openxmlformats.org/presentationml/2006/main' "
            b"xmlns:a='http://schemas.openxmlformats.org/drawingml/2006/main'>"
            b"<p:nvSpPr><p:cNvPr id='1' name='Shape'/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>"
            b"<p:spPr><a:xfrm><a:off x='0' y='0'/><a:ext cx='914400' cy='914400'/></a:xfrm>"
            b"<a:prstGeom prst='rect'/></p:spPr>"
            b"<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>"
            b"</p:sp>"
        ),
        None,
    )

    assert shape.text_frame._extents == (731520, 822960)


def test_font_properties_round_trip() -> None:
    font = Font(run(b"<a:rPr/>").get_or_add_rPr())

    assert font.bold is None
    font.bold = True
    assert font.bold is True

    assert font.italic is None
    font.italic = True
    assert font.italic is True

    assert font.language_id == MSO_LANGUAGE_ID.NONE
    font.language_id = MSO_LANGUAGE_ID.FRENCH
    assert font.language_id == MSO_LANGUAGE_ID.FRENCH
    font.language_id = MSO_LANGUAGE_ID.NONE
    assert font.language_id == MSO_LANGUAGE_ID.NONE

    assert font.name is None
    font.name = "Arial"
    assert font.name == "Arial"
    font.name = None
    assert font.name is None

    assert font.size is None
    font.size = Pt(12)
    assert font.size == Pt(12)
    font.size = None
    assert font.size is None

    assert font.underline is None
    font.underline = True
    assert font.underline is True
    font.underline = False
    assert font.underline is False
    font.underline = MSO_UNDERLINE.WAVY_LINE
    assert font.underline == MSO_UNDERLINE.WAVY_LINE


def test_font_color_and_fill() -> None:
    run_obj = _Run(run(b"<a:rPr/>"), None)
    color = run_obj.font.color

    color.rgb = RGBColor(0x12, 0x34, 0x56)
    assert color.rgb == RGBColor(0x12, 0x34, 0x56)

    color.theme_color = MSO_THEME_COLOR.ACCENT_1
    assert color.theme_color == MSO_THEME_COLOR.ACCENT_1

    color.brightness = 0.5
    assert color.brightness == 0.5


def test_paragraph_alignment_level_spacing() -> None:
    paragraph_obj = _Paragraph(paragraph(), None)

    assert paragraph_obj.alignment is None
    paragraph_obj.alignment = PP_ALIGN.CENTER
    assert paragraph_obj.alignment == PP_ALIGN.CENTER

    assert paragraph_obj.level == 0
    paragraph_obj.level = 2
    assert paragraph_obj.level == 2

    assert paragraph_obj.line_spacing is None
    paragraph_obj.line_spacing = 1.5
    assert paragraph_obj.line_spacing == 1.5
    paragraph_obj.line_spacing = Pt(18)
    assert paragraph_obj.line_spacing == Pt(18)

    assert paragraph_obj.space_before is None
    paragraph_obj.space_before = Pt(10)
    assert paragraph_obj.space_before == Pt(10)

    assert paragraph_obj.space_after is None
    paragraph_obj.space_after = Pt(12)
    assert paragraph_obj.space_after == Pt(12)


def test_paragraph_add_run_and_line_break(snapshot: SnapshotAssertion) -> None:
    paragraph_obj = _Paragraph(paragraph(), None)

    run_obj = paragraph_obj.add_run()
    run_obj.text = "Hello"
    paragraph_obj.add_line_break()

    assert serialize_xml(paragraph_obj._element) == snapshot


def test_paragraph_clear() -> None:
    paragraph_obj = _Paragraph(
        paragraph(b"<a:r><a:t>Hello</a:t></a:r><a:br/><a:r><a:t>World</a:t></a:r>"),
        None,
    )

    paragraph_obj.clear()

    assert paragraph_obj.text == ""
    assert paragraph_obj.runs == ()


def test_paragraph_text_setter_treats_newline_as_line_break() -> None:
    paragraph_obj = _Paragraph(paragraph(), None)

    paragraph_obj.text = "foo\nbar\vbaz"

    assert paragraph_obj.text == "foo\vbar\vbaz"


def test_run_text_escape_control_chars() -> None:
    run_obj = _Run(run(b"<a:rPr/><a:t/>"), None)

    run_obj.text = "a\x07b"

    assert "_x0007_" in run_obj._r.xml


def test_run_font_and_hyperlink_proxies(text_parent: DummyParent) -> None:
    run_obj = _Run(run(b"<a:rPr/><a:t>Hello</a:t>"), text_parent)

    assert isinstance(run_obj.font, Font)
    assert run_obj.hyperlink.address is None

    run_obj.hyperlink.address = "https://example.com"
    assert run_obj.hyperlink.address == "https://example.com"

    rel_ids = list(text_parent.part.rel_targets.keys())
    assert len(rel_ids) == 1

    run_obj.hyperlink.address = "https://example.org"
    assert run_obj.hyperlink.address == "https://example.org"
    assert len(text_parent.part.rel_targets) == 1

    run_obj.hyperlink.address = ""
    assert run_obj.hyperlink.address is None
    assert len(text_parent.part.rel_targets) == 0

    run_obj.hyperlink.address = "https://example.com"
    run_obj.hyperlink.address = None
    assert run_obj.hyperlink.address is None


def test_text_frame_part_property(text_parent: DummyParent) -> None:
    text_frame = TextFrame(tx_body(b"<a:bodyPr/><a:p/>"), text_parent)

    assert text_frame.part is text_parent.part
