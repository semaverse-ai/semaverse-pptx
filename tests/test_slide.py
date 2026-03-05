from __future__ import annotations

from dataclasses import dataclass, field

import pytest
from syrupy.assertion import SnapshotAssertion

from pptx.dml.fill import FillFormat
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml import parse_xml
from pptx.shapes.shapetree import (
    LayoutPlaceholders,
    LayoutShapes,
    MasterPlaceholders,
    MasterShapes,
    NotesSlidePlaceholders,
    NotesSlideShapes,
    SlidePlaceholders,
    SlideShapes,
)
from pptx.slide import (
    NotesMaster,
    NotesSlide,
    Slide,
    SlideLayout,
    SlideLayouts,
    SlideMaster,
    SlideMasters,
    Slides,
    _Background,
    _BaseMaster,
    _BaseSlide,
)
from pptx.text.text import TextFrame
from tests.factories import (
    notes_master_xml,
    notes_xml,
    slide_layout_xml,
    slide_master_xml,
    slide_xml,
)
from tests.stubs import (
    CloneRecorder,
    NewSlideStub,
    ParentProxy,
    SlidePartStub,
    SlidesPartStub,
)
from tests.xml_utils import serialize_xml


def test_base_slide_name_get_set_and_clear(snapshot: SnapshotAssertion) -> None:
    # Arrange
    element = parse_xml(
        b"""
        <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:cSld name="Slide 1"/>
        </p:sld>
        """
    )
    base_slide = _BaseSlide(element, None)  # type: ignore[arg-type]

    # Act
    base_slide.name = "Agenda"
    base_slide.name = None

    # Assert
    assert base_slide.name == ""
    assert serialize_xml(element) == snapshot


def test_base_slide_background_property_returns_background_proxy() -> None:
    # Arrange
    slide = _BaseSlide(parse_xml(slide_xml()), None)  # type: ignore[arg-type]

    # Act
    background = slide.background

    # Assert
    assert isinstance(background, _Background)


def test_base_master_provides_shapes_and_placeholders() -> None:
    # Arrange
    master_xml = slide_master_xml()
    base_master = _BaseMaster(parse_xml(master_xml), None)  # type: ignore[arg-type]

    # Act
    placeholders = base_master.placeholders
    shapes = base_master.shapes

    # Assert
    assert isinstance(placeholders, MasterPlaceholders)
    assert isinstance(shapes, MasterShapes)


def test_notes_slide_provides_shapes_and_placeholders() -> None:
    # Arrange
    notes_slide = NotesSlide(parse_xml(notes_xml()), None)  # type: ignore[arg-type]

    # Act
    placeholders = notes_slide.placeholders
    shapes = notes_slide.shapes

    # Assert
    assert isinstance(placeholders, NotesSlidePlaceholders)
    assert isinstance(shapes, NotesSlideShapes)


def test_notes_slide_notes_placeholder_and_notes_text_frame() -> None:
    # Arrange
    body_placeholder = b"""
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Notes Placeholder"/>
          <p:cNvSpPr/>
          <p:nvPr><p:ph type="body"/></p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
        <p:txBody>
          <a:bodyPr/>
          <a:lstStyle/>
          <a:p/>
        </p:txBody>
      </p:sp>
    """
    notes_slide = NotesSlide(parse_xml(notes_xml(body_placeholder)), None)  # type: ignore[arg-type]

    # Act
    notes_placeholder = notes_slide.notes_placeholder
    notes_text_frame = notes_slide.notes_text_frame

    # Assert
    assert notes_placeholder is not None
    assert notes_placeholder.placeholder_format.type == PP_PLACEHOLDER.BODY
    assert isinstance(notes_text_frame, TextFrame)


def test_notes_slide_notes_text_frame_is_none_when_body_placeholder_missing() -> None:
    # Arrange
    notes_slide = NotesSlide(parse_xml(notes_xml()), None)  # type: ignore[arg-type]

    # Act / Assert
    assert notes_slide.notes_placeholder is None
    assert notes_slide.notes_text_frame is None


def test_notes_slide_clone_master_placeholders_copies_only_cloneable_types() -> None:
    # Arrange
    master_shapes = b"""
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Slide Image"/>
          <p:cNvSpPr/>
          <p:nvPr><p:ph type="sldImg"/></p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
      </p:sp>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="3" name="Body"/>
          <p:cNvSpPr/>
          <p:nvPr><p:ph type="body"/></p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
      </p:sp>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="4" name="Header"/>
          <p:cNvSpPr/>
          <p:nvPr><p:ph type="hdr"/></p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
      </p:sp>
    """
    notes_master = NotesMaster(parse_xml(notes_master_xml(master_shapes)), None)  # type: ignore[arg-type]
    notes_slide = NotesSlide(parse_xml(notes_xml()), None)  # type: ignore[arg-type]

    # Act
    notes_slide.clone_master_placeholders(notes_master)
    cloned_types = [placeholder.placeholder_format.type for placeholder in notes_slide.placeholders]

    # Assert
    assert PP_PLACEHOLDER.SLIDE_IMAGE in cloned_types
    assert PP_PLACEHOLDER.BODY in cloned_types
    assert PP_PLACEHOLDER.HEADER not in cloned_types


@pytest.mark.parametrize(("with_bg", "expected"), [(False, True), (True, False)])
def test_slide_follow_master_background(with_bg: bool, expected: bool) -> None:
    # Arrange
    slide = Slide(parse_xml(slide_xml(with_bg=with_bg)), SlidePartStub())  # type: ignore[arg-type]

    # Act / Assert
    assert slide.follow_master_background is expected


def test_slide_delegates_part_backed_properties() -> None:
    # Arrange
    part = SlidePartStub()
    slide = Slide(parse_xml(slide_xml()), part)  # type: ignore[arg-type]

    # Act / Assert
    assert slide.has_notes_slide is True
    assert slide.slide_id == 256
    assert slide.slide_layout == "layout"
    assert slide.notes_slide == "notes"
    assert isinstance(slide.placeholders, SlidePlaceholders)
    assert isinstance(slide.shapes, SlideShapes)


def test_slides_collection_semantics() -> None:
    # Arrange
    sld_id_lst = parse_xml(
        b"""
        <p:sldIdLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                    xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <p:sldId id="256" r:id="rId1"/>
          <p:sldId id="257" r:id="rId2"/>
        </p:sldIdLst>
        """
    )
    slide_1 = object()
    slide_2 = object()
    part = SlidesPartStub(
        slides_by_rid={"rId1": slide_1, "rId2": slide_2},
        slide_by_id={256: slide_1, 257: slide_2},
        add_slide_result=("unused", object()),
    )
    slides = Slides(sld_id_lst, ParentProxy(part=part))  # type: ignore[arg-type]

    # Act / Assert
    assert len(slides) == 2
    assert slides[0] is slide_1
    assert list(slides) == [slide_1, slide_2]
    assert slides.get(257) is slide_2
    assert slides.get(999, default="missing") == "missing"
    assert slides.index(slide_2) == 1

    with pytest.raises(ValueError):
        slides.index(object())

    with pytest.raises(IndexError, match="slide index out of range"):
        _ = slides[2]


def test_slides_add_slide_mutates_sldidlst_and_clones_layout_placeholders(
    snapshot: SnapshotAssertion,
) -> None:
    # Arrange
    sld_id_lst = parse_xml(
        b"""
        <p:sldIdLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                    xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <p:sldId id="256" r:id="rId1"/>
        </p:sldIdLst>
        """
    )
    recorder = CloneRecorder()
    new_slide = NewSlideStub(shapes=recorder)
    part = SlidesPartStub(
        slides_by_rid={},
        slide_by_id={},
        add_slide_result=("rId2", new_slide),
    )
    slides = Slides(sld_id_lst, ParentProxy(part=part))  # type: ignore[arg-type]
    layout = object()

    # Act
    result = slides.add_slide(layout)  # type: ignore[arg-type]

    # Assert
    assert result is new_slide
    assert recorder.cloned_with == [layout]
    assert serialize_xml(sld_id_lst) == snapshot


def test_slide_layout_cloneable_placeholders_filters_latent_types() -> None:
    # Arrange
    layout_shapes = b"""
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Date Placeholder"/>
          <p:cNvSpPr/>
          <p:nvPr><p:ph type="dt"/></p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
      </p:sp>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="3" name="Body Placeholder"/>
          <p:cNvSpPr/>
          <p:nvPr><p:ph type="body"/></p:nvPr>
        </p:nvSpPr>
        <p:spPr/>
      </p:sp>
    """
    layout = SlideLayout(parse_xml(slide_layout_xml(layout_shapes)), _LayoutPartStub())  # type: ignore[arg-type]

    # Act
    cloneable = list(layout.iter_cloneable_placeholders())
    cloneable_types = [placeholder.placeholder_format.type for placeholder in cloneable]

    # Assert
    assert cloneable_types == [PP_PLACEHOLDER.BODY]


@dataclass
class _LayoutPartStub:
    slide_master: object = "slide-master"
    package: object | None = None


def test_slide_layout_properties_and_used_by_slides() -> None:
    # Arrange
    part = _LayoutPartStub()
    layout = SlideLayout(parse_xml(slide_layout_xml()), part)  # type: ignore[arg-type]

    class _SlideStub:
        def __init__(self, slide_layout: object) -> None:
            self.slide_layout = slide_layout

    included_slide = _SlideStub(layout)
    excluded_slide = _SlideStub("other-layout")

    @dataclass
    class _PresentationStub:
        slides: list[object]

    @dataclass
    class _PresentationPartStub:
        presentation: _PresentationStub

    @dataclass
    class _PackageStub:
        presentation_part: _PresentationPartStub

    part.package = _PackageStub(
        _PresentationPartStub(_PresentationStub([included_slide, excluded_slide]))
    )

    # Act
    used_by_slides = layout.used_by_slides

    # Assert
    assert isinstance(layout.placeholders, LayoutPlaceholders)
    assert isinstance(layout.shapes, LayoutShapes)
    assert layout.slide_master == "slide-master"
    assert used_by_slides == (included_slide,)


@dataclass
class _DropRelPartStub:
    dropped: list[str] = field(default_factory=list)

    def drop_rel(self, rid: str) -> None:
        self.dropped.append(rid)


@dataclass
class _SlideMasterProxy:
    part: _DropRelPartStub


@dataclass
class _SlideLayoutStub:
    name: str
    used_by_slides: tuple[object, ...]
    slide_master: _SlideMasterProxy


@dataclass
class _SlideLayoutsPartStub:
    layouts_by_rid: dict[str, _SlideLayoutStub]

    def related_slide_layout(self, rid: str) -> _SlideLayoutStub:
        return self.layouts_by_rid[rid]


def test_slide_layouts_collection_semantics_and_remove() -> None:
    # Arrange
    sld_layout_id_lst = parse_xml(
        b"""
        <p:sldLayoutIdLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <p:sldLayoutId r:id="rId1"/>
          <p:sldLayoutId r:id="rId2"/>
        </p:sldLayoutIdLst>
        """
    )
    drop_rel_part = _DropRelPartStub()
    free_layout = _SlideLayoutStub(
        name="Title Slide",
        used_by_slides=(),
        slide_master=_SlideMasterProxy(part=drop_rel_part),
    )
    in_use_layout = _SlideLayoutStub(
        name="Blank",
        used_by_slides=(object(),),
        slide_master=_SlideMasterProxy(part=drop_rel_part),
    )
    part = _SlideLayoutsPartStub(layouts_by_rid={"rId1": free_layout, "rId2": in_use_layout})
    layouts = SlideLayouts(sld_layout_id_lst, ParentProxy(part=part))  # type: ignore[arg-type]

    # Act / Assert
    assert len(layouts) == 2
    assert layouts[0] is free_layout
    assert list(layouts) == [free_layout, in_use_layout]
    assert layouts.get_by_name("Blank") is in_use_layout
    assert layouts.get_by_name("Missing", default=None) is None
    assert layouts.index(in_use_layout) == 1

    with pytest.raises(IndexError, match="slide layout index out of range"):
        _ = layouts[2]

    with pytest.raises(ValueError, match="cannot remove slide-layout in use"):
        layouts.remove(in_use_layout)  # type: ignore[arg-type]

    layouts.remove(free_layout)  # type: ignore[arg-type]
    assert len(layouts) == 1
    assert sld_layout_id_lst.sldLayoutId_lst[0].rId == "rId2"
    assert drop_rel_part.dropped == ["rId1"]


@dataclass
class _SlideMastersPartStub:
    masters_by_rid: dict[str, object]

    def related_slide_master(self, rid: str) -> object:
        return self.masters_by_rid[rid]


def test_slide_master_and_slide_masters_collection() -> None:
    # Arrange
    master_xml = slide_master_xml()
    slide_master = SlideMaster(parse_xml(master_xml), _LayoutPartStub())  # type: ignore[arg-type]

    masters_xml = parse_xml(
        b"""
        <p:sldMasterIdLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <p:sldMasterId r:id="rId1"/>
        </p:sldMasterIdLst>
        """
    )
    part = _SlideMastersPartStub({"rId1": slide_master})
    masters = SlideMasters(masters_xml, ParentProxy(part=part))  # type: ignore[arg-type]

    # Act / Assert
    assert isinstance(slide_master.slide_layouts, SlideLayouts)
    assert len(masters) == 1
    assert masters[0] is slide_master
    assert list(masters) == [slide_master]

    with pytest.raises(IndexError, match="slide master index out of range"):
        _ = masters[1]


def test_background_fill_returns_fill_format() -> None:
    # Arrange
    c_sld = parse_xml(
        b"""
        <p:cSld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <p:bg>
            <p:bgPr>
              <a:solidFill/>
            </p:bgPr>
          </p:bg>
        </p:cSld>
        """
    )
    background = _Background(c_sld)

    # Act
    fill = background.fill

    # Assert
    assert isinstance(fill, FillFormat)
