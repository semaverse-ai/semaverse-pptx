from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.freeform import FreeformBuilder, _Close, _LineSegment, _MoveTo
from pptx.shapes.shapetree import SlideShapes


def test_freeform_builder_new(parent) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)

    builder = FreeformBuilder.new(shapes, 10.4, 20.6, 2.0, 3.0)

    assert isinstance(builder, FreeformBuilder)
    assert builder._start_x == 10
    assert builder._start_y == 21


def test_freeform_builder_add_line_segments_and_move_to(parent) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)
    builder = FreeformBuilder.new(shapes, 0, 0, 1.0, 1.0)

    result = builder.add_line_segments([(10, 20), (30, 40)], close=False).move_to(50, 60)

    assert result is builder
    assert len(builder) == 3
    assert isinstance(builder[0], _LineSegment)
    assert isinstance(builder[1], _LineSegment)
    assert isinstance(builder[2], _MoveTo)


def test_freeform_builder_convert_to_shape(parent) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)
    builder = FreeformBuilder.new(shapes, 10, 20, 1.0, 1.0)
    builder.add_line_segments([(30, 40), (50, 60)])
    builder.move_to(100, 100)
    builder.add_line_segments([(150, 150)], close=False)

    shape = builder.convert_to_shape(origin_x=5, origin_y=5)

    assert shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
    assert shape.left == 15
    assert shape.top == 25
    assert shape.width == 140
    assert shape.height == 130
    sp = parent.part._element.cSld.spTree[-1]
    assert sp.tag.endswith("sp")
    assert sp.spPr.custGeom is not None
    assert sp.spPr.custGeom.pathLst is not None


def test_freeform_builder_geometry_properties(parent) -> None:
    builder = FreeformBuilder.new(
        SlideShapes(parent.part._element.cSld.spTree, parent), 10, 20, 2.0, 3.0
    )
    builder.add_line_segments([(20, 30), (40, 50)], close=False)

    assert builder.shape_offset_x == 10
    assert builder.shape_offset_y == 20
    assert builder._dx == 30
    assert builder._dy == 30
    assert builder._width == 60
    assert builder._height == 90


def test_drawing_operations_apply_to_path(parent) -> None:
    shapes = SlideShapes(parent.part._element.cSld.spTree, parent)
    builder = FreeformBuilder.new(shapes, 0, 0, 1.0, 1.0)
    sp = builder._add_freeform_sp(0, 0)
    path = builder._start_path(sp)

    _LineSegment.new(builder, 25, 35).apply_operation_to(path)
    _MoveTo.new(builder, 45, 55).apply_operation_to(path)
    _Close.new().apply_operation_to(path)

    assert len(path) == 4
    assert path[0].tag.endswith("moveTo")
    assert path[1].tag.endswith("lnTo")
    assert path[2].tag.endswith("moveTo")
    assert path[3].tag.endswith("close")
