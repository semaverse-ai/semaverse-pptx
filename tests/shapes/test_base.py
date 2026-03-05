from __future__ import annotations

import pytest

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape, _PlaceholderFormat

from .factories import make_sp


def test_base_shape_properties(parent) -> None:
    shape = BaseShape(
        make_sp(shape_id=42, name="Test Shape", x=100, y=200, cx=300, cy=400, rot=60000),
        parent,
    )

    assert shape.shape_id == 42
    assert shape.name == "Test Shape"
    assert shape.left == 100
    assert shape.top == 200
    assert shape.width == 300
    assert shape.height == 400
    assert shape.rotation == 1.0
    assert shape.part is parent.part
    assert shape.has_chart is False
    assert shape.has_table is False
    assert shape.has_text_frame is False
    assert shape.is_placeholder is False


def test_base_shape_setters(parent) -> None:
    shape = BaseShape(
        make_sp(shape_id=42, name="Old Name", x=100, y=200, cx=300, cy=400),
        parent,
    )

    shape.name = "New Name"
    shape.left = 110
    shape.top = 220
    shape.width = 330
    shape.height = 440
    shape.rotation = 45.0

    assert shape.name == "New Name"
    assert shape.left == 110
    assert shape.top == 220
    assert shape.width == 330
    assert shape.height == 440
    assert shape.rotation == 45.0


def test_base_shape_equality(parent) -> None:
    element = parse_xml(
        b'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
    )

    one = BaseShape(element, parent)
    two = BaseShape(element, parent)
    three = BaseShape(
        parse_xml(b'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'),
        parent,
    )

    assert one == two
    assert one != three
    assert one != "not-shape"


def test_placeholder_format(parent) -> None:
    shape = BaseShape(
        make_sp(shape_id=1, name="Title", placeholder_attrs='type="title" idx="2"'),
        parent,
    )

    placeholder_format = shape.placeholder_format

    assert isinstance(placeholder_format, _PlaceholderFormat)
    assert placeholder_format.idx == 2
    assert placeholder_format.type == PP_PLACEHOLDER.TITLE


def test_placeholder_format_raises_for_non_placeholder(parent) -> None:
    shape = BaseShape(
        make_sp(shape_id=1, name="No Placeholder"),
        parent,
    )

    with pytest.raises(ValueError, match="shape is not a placeholder"):
        _ = shape.placeholder_format
